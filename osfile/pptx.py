from pptx.util import Inches, Pt
from PIL import Image


def pptx_slide_with_image(prs, image_path, title=None, font_name="Arial", font_size=24):
    # Taille de la diapositive (par défaut 10 x 7.5 pouces)
    slide_width_inch = prs.slide_width / Inches(1)
    slide_height_inch = prs.slide_height / Inches(1)

    if title:
        slide_layout = prs.slide_layouts[5]  # Utiliser une mise en page avec un titre
        slide = prs.slides.add_slide(slide_layout)
        # Ajouter le titre avec le style personnalisé
        title_shape = slide.shapes.title
        title_shape.text = title
        text_frame = title_shape.text_frame
        p = text_frame.paragraphs[0]

        # Modifier le style de police du titre
        font = p.runs[0].font  # Le premier run contient déjà le texte
        font.name = font_name  # Nom de la police (par exemple, 'Arial')
        font.size = Pt(font_size)  # Taille de la police en points
    else:
        slide_layout = prs.slide_layouts[6]  # Utiliser une mise en page avec un titre
        slide = prs.slides.add_slide(slide_layout)

    # Charger l'image pour obtenir ses dimensions avec Pillow
    img = Image.open(image_path)
    img_width_px, img_height_px = img.size

    # Convertir les dimensions en pouces (pptx utilise une résolution de 72 DPI par défaut)
    img_width_inch = img_width_px / 72
    img_height_inch = img_height_px / 72

    # Calculer les positions pour centrer l'image
    if img_width_inch > slide_width_inch or img_height_inch > slide_height_inch:
        # Redimensionner si l'image est plus grande que la diapositive
        scale_factor = min(
            slide_width_inch / img_width_inch, slide_height_inch / img_height_inch
        )
        img_width_inch *= scale_factor
        img_height_inch *= scale_factor

    # Calculer les positions pour centrer l'image
    left = (Inches(slide_width_inch) - Inches(img_width_inch)) / 2
    top = (Inches(slide_height_inch) - Inches(img_height_inch)) / 2

    # Ajouter l'image centrée et redimensionnée automatiquement
    slide.shapes.add_picture(
        image_path,
        left,
        top,
        width=Inches(img_width_inch),
        height=Inches(img_height_inch),
    )